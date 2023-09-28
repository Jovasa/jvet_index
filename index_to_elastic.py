import subprocess
from os import environ
from pathlib import Path

import elasticsearch as es
import PyPDF2
import textract
from docx import Document
from pptx import Presentation

environ['ES_CERT_FILE'] = "/home/jovasa/PycharmProjects/IndexJVET/http_cs.crt"
environ['ES_API_ID'] = "wCnfo4oBtbHqigGW_LMp"
environ['ES_API_KEY'] = "grbWNaE9ReKSU7p75zrZSw"

client = es.Elasticsearch([{'host': 'localhost', 'port': 9500, 'scheme': 'https'}],
                          basic_auth=("elastic", "*iiFGM_ss*YOGU0x1pOC"),
                          verify_certs=False,
                          ssl_show_warn=False)


# client.indices.create(index="jvet_documents")


def index_one_zip_file(file_path: Path):
    base_out = file_path.with_suffix("")
    base_out.mkdir(exist_ok=True)
    try:
        subprocess.check_call(["unzip", "-d", str(base_out), str(file_path)], stdout=subprocess.DEVNULL)
    except subprocess.CalledProcessError:
        print("Error unzipping", file_path)
        # base_out.rmdir()
        return

    file_type_counts = {
        ".doc": 0,
        ".docx": 0,
        ".pdf": 0,
        ".pptx": 0,
    }
    files = [x for x in base_out.rglob("*") if x.is_file()]
    for file in files:
        try:
            if file.suffix == ".doc":
                d = textract.process(str(file)).decode("utf-8")
                # print(a.decode("utf-8"))
            elif file.suffix == ".docx":
                doc = Document(file)
                d = []
                for p in doc.paragraphs:
                    d.extend(p.text.split("\n"))
                d = "\n".join(d)
                # print(*d,sep="\n")
            elif file.suffix == ".pdf":
                reader = PyPDF2.PdfReader(file)
                d = []
                for page in reader.pages:
                    d.extend(page.extract_text())
                d = "\n".join(d)
            elif file.suffix == ".pptx":
                p = Presentation(file)
                d = []
                for slide in p.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            d.append(shape.text)
                d = "\n".join(d)
            else:
                continue
        except Exception as e:
            print("Error processing", file)
            continue
        file_type_counts[file.suffix] += 1
        client.index(index="jvet_documents",

                     id=file_path.stem + file.suffix + str(file_type_counts[file.suffix]),
                     document={
                         "file_name": file.name,
                         "file_type": file.suffix,
                         "text": d,
                     })

    # remove the files
    try:
        subprocess.check_call(["rm", "-rf", str(base_out)])
    except subprocess.CalledProcessError:
        print("Error removing", base_out)


if __name__ == '__main__':
    all_zip_files = [x for x in Path("contents").iterdir() if x.is_file()]
    for fip_file in all_zip_files:
        index_one_zip_file(fip_file)
